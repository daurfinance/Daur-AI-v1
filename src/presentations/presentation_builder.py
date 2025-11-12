#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
Daur-AI: Модуль создания профессиональных презентаций
Включает создание, редактирование и экспорт презентаций в различные форматы

Версия: 2.0
Дата: 25.10.2025
Автор: Manus AI
"""

import logging
import json
import os
from typing import Dict, List, Any, Optional, Tuple
from datetime import datetime
from pathlib import Path
from enum import Enum
from dataclasses import dataclass, asdict
import xml.etree.ElementTree as ET


class SlideLayout(Enum):
    """Типы макетов слайдов"""
    TITLE = "title"
    TITLE_CONTENT = "title_content"
    TWO_COLUMN = "two_column"
    BLANK = "blank"
    SECTION = "section"
    CLOSING = "closing"


class TransitionType(Enum):
    """Типы переходов между слайдами"""
    NONE = "none"
    FADE = "fade"
    PUSH = "push"
    WIPE = "wipe"
    COVER = "cover"
    UNCOVER = "uncover"
    SPLIT = "split"
    ZOOM = "zoom"


@dataclass
class TextStyle:
    """Стиль текста"""
    font_name: str = "Calibri"
    font_size: int = 18
    bold: bool = False
    italic: bool = False
    underline: bool = False
    color: str = "#000000"
    alignment: str = "left"  # left, center, right


@dataclass
class SlideContent:
    """Содержимое слайда"""
    title: str = ""
    subtitle: str = ""
    content: List[str] = None
    images: List[str] = None
    notes: str = ""
    layout: SlideLayout = SlideLayout.TITLE_CONTENT
    background_color: str = "#FFFFFF"
    transition: TransitionType = TransitionType.FADE
    transition_duration: float = 0.5
    
    def __post_init__(self):
        if self.content is None:
            self.content = []
        if self.images is None:
            self.images = []


class PresentationTheme:
    """Тема презентации"""
    
    def __init__(self, name: str = "default"):
        """
        Args:
            name: Имя темы
        """
        self.name = name
        self.colors = {
            'primary': '#1F77B4',
            'secondary': '#FF7F0E',
            'accent': '#2CA02C',
            'background': '#FFFFFF',
            'text': '#000000',
            'text_light': '#666666'
        }
        self.fonts = {
            'title': 'Arial',
            'content': 'Calibri',
            'code': 'Courier New'
        }
        self.sizes = {
            'title': 44,
            'subtitle': 32,
            'heading': 28,
            'content': 18,
            'small': 14
        }
    
    def set_color(self, color_name: str, hex_color: str):
        """Установить цвет темы"""
        self.colors[color_name] = hex_color
    
    def set_font(self, font_type: str, font_name: str):
        """Установить шрифт"""
        self.fonts[font_type] = font_name
    
    def set_size(self, size_type: str, size: int):
        """Установить размер"""
        self.sizes[size_type] = size
    
    def to_dict(self) -> Dict[str, Any]:
        """Преобразовать в словарь"""
        return {
            'name': self.name,
            'colors': self.colors,
            'fonts': self.fonts,
            'sizes': self.sizes
        }


class Presentation:
    """Основной класс презентации"""
    
    def __init__(self, title: str, author: str = "Daur AI", theme: Optional[PresentationTheme] = None):
        """
        Args:
            title: Название презентации
            author: Автор
            theme: Тема презентации
        """
        self.title = title
        self.author = author
        self.theme = theme or PresentationTheme()
        self.slides: List[SlideContent] = []
        self.created_at = datetime.now()
        self.modified_at = datetime.now()
        self.logger = logging.getLogger('daur_ai.presentation_builder')
    
    def add_slide(self, slide: SlideContent) -> int:
        """
        Добавить слайд
        
        Args:
            slide: Содержимое слайда
            
        Returns:
            int: Индекс слайда
        """
        self.slides.append(slide)
        self.modified_at = datetime.now()
        self.logger.info(f"Слайд добавлен: {slide.title}")
        return len(self.slides) - 1
    
    def add_title_slide(self, title: str, subtitle: str = "") -> int:
        """Добавить титульный слайд"""
        slide = SlideContent(
            title=title,
            subtitle=subtitle,
            layout=SlideLayout.TITLE
        )
        return self.add_slide(slide)
    
    def add_content_slide(self, title: str, content: List[str], images: List[str] = None) -> int:
        """Добавить слайд с содержимым"""
        slide = SlideContent(
            title=title,
            content=content or [],
            images=images or [],
            layout=SlideLayout.TITLE_CONTENT
        )
        return self.add_slide(slide)
    
    def add_section_slide(self, title: str, subtitle: str = "") -> int:
        """Добавить слайд раздела"""
        slide = SlideContent(
            title=title,
            subtitle=subtitle,
            layout=SlideLayout.SECTION
        )
        return self.add_slide(slide)
    
    def add_closing_slide(self, title: str = "Спасибо!", subtitle: str = "") -> int:
        """Добавить закрывающий слайд"""
        slide = SlideContent(
            title=title,
            subtitle=subtitle,
            layout=SlideLayout.CLOSING
        )
        return self.add_slide(slide)
    
    def edit_slide(self, index: int, slide: SlideContent):
        """
        Редактировать слайд
        
        Args:
            index: Индекс слайда
            slide: Новое содержимое
        """
        if 0 <= index < len(self.slides):
            self.slides[index] = slide
            self.modified_at = datetime.now()
            self.logger.info(f"Слайд {index} отредактирован")
    
    def delete_slide(self, index: int):
        """Удалить слайд"""
        if 0 <= index < len(self.slides):
            del self.slides[index]
            self.modified_at = datetime.now()
            self.logger.info(f"Слайд {index} удален")
    
    def move_slide(self, from_index: int, to_index: int):
        """Переместить слайд"""
        if 0 <= from_index < len(self.slides) and 0 <= to_index < len(self.slides):
            slide = self.slides.pop(from_index)
            self.slides.insert(to_index, slide)
            self.modified_at = datetime.now()
    
    def get_slide_count(self) -> int:
        """Получить количество слайдов"""
        return len(self.slides)
    
    def get_slide(self, index: int) -> Optional[SlideContent]:
        """Получить слайд"""
        if 0 <= index < len(self.slides):
            return self.slides[index]
        return None
    
    def get_all_slides(self) -> List[SlideContent]:
        """Получить все слайды"""
        return self.slides.copy()
    
    def to_dict(self) -> Dict[str, Any]:
        """Преобразовать в словарь"""
        return {
            'title': self.title,
            'author': self.author,
            'created_at': self.created_at.isoformat(),
            'modified_at': self.modified_at.isoformat(),
            'theme': self.theme.to_dict(),
            'slides': [
                {
                    'title': slide.title,
                    'subtitle': slide.subtitle,
                    'content': slide.content,
                    'images': slide.images,
                    'notes': slide.notes,
                    'layout': slide.layout.value,
                    'background_color': slide.background_color,
                    'transition': slide.transition.value,
                    'transition_duration': slide.transition_duration
                }
                for slide in self.slides
            ]
        }
    
    def save_json(self, filepath: str):
        """Сохранить презентацию в JSON"""
        with open(filepath, 'w', encoding='utf-8') as f:
            json.dump(self.to_dict(), f, indent=2, ensure_ascii=False)
        self.logger.info(f"Презентация сохранена в {filepath}")
    
    @classmethod
    def load_json(cls, filepath: str) -> 'Presentation':
        """Загрузить презентацию из JSON"""
        with open(filepath, 'r', encoding='utf-8') as f:
            data = json.load(f)
        
        # Создаем презентацию
        pres = cls(data['title'], data['author'])
        
        # Загружаем слайды
        for slide_data in data.get('slides', []):
            slide = SlideContent(
                title=slide_data.get('title', ''),
                subtitle=slide_data.get('subtitle', ''),
                content=slide_data.get('content', []),
                images=slide_data.get('images', []),
                notes=slide_data.get('notes', ''),
                layout=SlideLayout(slide_data.get('layout', 'title_content')),
                background_color=slide_data.get('background_color', '#FFFFFF'),
                transition=TransitionType(slide_data.get('transition', 'fade')),
                transition_duration=slide_data.get('transition_duration', 0.5)
            )
            pres.add_slide(slide)
        
        return pres


class PresentationExporter:
    """Экспортер презентаций"""
    
    def __init__(self):
        """Инициализация"""
        self.logger = logging.getLogger('daur_ai.presentation_exporter')
    
    def export_to_pptx(self, presentation: Presentation, filepath: str):
        """
        Экспортировать в PPTX
        
        Args:
            presentation: Презентация
            filepath: Путь к файлу
        """
        try:
            from pptx import Presentation as PPTXPresentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor
            
            # Создаем PPTX презентацию
            prs = PPTXPresentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # Добавляем слайды
            for slide_content in presentation.slides:
                # Выбираем макет
                if slide_content.layout == SlideLayout.TITLE:
                    layout = prs.slide_layouts[0]  # Title slide
                elif slide_content.layout == SlideLayout.TITLE_CONTENT:
                    layout = prs.slide_layouts[1]  # Title and content
                else:
                    layout = prs.slide_layouts[6]  # Blank
                
                slide = prs.slides.add_slide(layout)
                
                # Добавляем заголовок
                if slide_content.title and len(slide.shapes) > 0:
                    title_shape = slide.shapes.title
                    title_shape.text = slide_content.title
                
                # Добавляем содержимое
                if slide_content.content and len(slide.shapes) > 1:
                    body_shape = slide.placeholders[1]
                    text_frame = body_shape.text_frame
                    text_frame.clear()
                    
                    for item in slide_content.content:
                        p = text_frame.add_paragraph()
                        p.text = item
                        p.level = 0
            
            # Сохраняем
            prs.save(filepath)
            self.logger.info(f"Презентация экспортирована в PPTX: {filepath}")
        
        except ImportError:
            self.logger.error("python-pptx не установлен. Установите: pip install python-pptx")
    
    def export_to_pdf(self, presentation: Presentation, filepath: str):
        """
        Экспортировать в PDF
        
        Args:
            presentation: Презентация
            filepath: Путь к файлу
        """
        try:
            from reportlab.lib.pagesizes import landscape, letter
            from reportlab.pdfgen import canvas
            from reportlab.lib.units import inch
            
            # Создаем PDF
            c = canvas.Canvas(filepath, pagesize=landscape(letter))
            width, height = landscape(letter)
            
            # Добавляем слайды
            for i, slide in enumerate(presentation.slides):
                # Фон
                c.setFillColor(slide.background_color)
                c.rect(0, 0, width, height, fill=1)
                
                # Заголовок
                if slide.title:
                    c.setFont(presentation.theme.fonts['title'], presentation.theme.sizes['title'])
                    c.drawString(0.5 * inch, height - 0.75 * inch, slide.title)
                
                # Содержимое
                y = height - 1.5 * inch
                for item in slide.content:
                    c.setFont(presentation.theme.fonts['content'], presentation.theme.sizes['content'])
                    c.drawString(0.75 * inch, y, item)
                    y -= 0.4 * inch
                
                c.showPage()
            
            c.save()
            self.logger.info(f"Презентация экспортирована в PDF: {filepath}")
        
        except ImportError:
            self.logger.error("reportlab не установлен. Установите: pip install reportlab")
    
    def export_to_html(self, presentation: Presentation, filepath: str):
        """
        Экспортировать в HTML
        
        Args:
            presentation: Презентация
            filepath: Путь к файлу
        """
        html_content = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>{presentation.title}</title>
    <style>
        body {{
            font-family: {presentation.theme.fonts['content']};
            margin: 0;
            padding: 0;
            background-color: #f0f0f0;
        }}
        .slide {{
            width: 100%;
            height: 100vh;
            background-color: {presentation.theme.colors['background']};
            color: {presentation.theme.colors['text']};
            padding: 40px;
            box-sizing: border-box;
            page-break-after: always;
            display: flex;
            flex-direction: column;
            justify-content: center;
        }}
        .slide h1 {{
            font-size: {presentation.theme.sizes['title']}px;
            color: {presentation.theme.colors['primary']};
            margin: 0 0 20px 0;
        }}
        .slide h2 {{
            font-size: {presentation.theme.sizes['subtitle']}px;
            color: {presentation.theme.colors['secondary']};
            margin: 0 0 20px 0;
        }}
        .slide p {{
            font-size: {presentation.theme.sizes['content']}px;
            margin: 10px 0;
        }}
        .slide img {{
            max-width: 100%;
            max-height: 400px;
            margin: 20px 0;
        }}
    </style>
</head>
<body>
"""
        
        # Добавляем слайды
        for slide in presentation.slides:
            html_content += f"""    <div class="slide" style="background-color: {slide.background_color};">
"""
            
            if slide.title:
                html_content += f"""        <h1>{slide.title}</h1>
"""
            
            if slide.subtitle:
                html_content += f"""        <h2>{slide.subtitle}</h2>
"""
            
            for item in slide.content:
                html_content += f"""        <p>{item}</p>
"""
            
            for image in slide.images:
                html_content += f"""        <img src="{image}" alt="Image">
"""
            
            html_content += """    </div>
"""
        
        html_content += """</body>
</html>"""
        
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(html_content)
        
        self.logger.info(f"Презентация экспортирована в HTML: {filepath}")


class PresentationBuilder:
    """Построитель презентаций с удобным API"""
    
    def __init__(self, title: str):
        """
        Args:
            title: Название презентации
        """
        self.presentation = Presentation(title)
        self.logger = logging.getLogger('daur_ai.presentation_builder')
    
    def set_theme(self, theme: PresentationTheme) -> 'PresentationBuilder':
        """Установить тему"""
        self.presentation.theme = theme
        return self
    
    def add_title(self, title: str, subtitle: str = "") -> 'PresentationBuilder':
        """Добавить титульный слайд"""
        self.presentation.add_title_slide(title, subtitle)
        return self
    
    def add_section(self, title: str, subtitle: str = "") -> 'PresentationBuilder':
        """Добавить слайд раздела"""
        self.presentation.add_section_slide(title, subtitle)
        return self
    
    def add_content(self, title: str, *items) -> 'PresentationBuilder':
        """Добавить слайд с содержимым"""
        self.presentation.add_content_slide(title, list(items))
        return self
    
    def add_closing(self, title: str = "Спасибо!") -> 'PresentationBuilder':
        """Добавить закрывающий слайд"""
        self.presentation.add_closing_slide(title)
        return self
    
    def build(self) -> Presentation:
        """Построить презентацию"""
        return self.presentation
    
    def save(self, filepath: str):
        """Сохранить презентацию"""
        self.presentation.save_json(filepath)
        return self
    
    def export_pptx(self, filepath: str):
        """Экспортировать в PPTX"""
        exporter = PresentationExporter()
        exporter.export_to_pptx(self.presentation, filepath)
        return self
    
    def export_pdf(self, filepath: str):
        """Экспортировать в PDF"""
        exporter = PresentationExporter()
        exporter.export_to_pdf(self.presentation, filepath)
        return self
    
    def export_html(self, filepath: str):
        """Экспортировать в HTML"""
        exporter = PresentationExporter()
        exporter.export_to_html(self.presentation, filepath)
        return self


# Глобальный экземпляр
_presentation_builder = None


def get_presentation_builder(title: str) -> PresentationBuilder:
    """Получить построитель презентаций"""
    return PresentationBuilder(title)

